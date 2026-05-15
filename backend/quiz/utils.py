import os
import json
import re
import ast
import random
from groq import Groq
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

client = Groq(api_key=os.getenv('GROQ_API_KEY'))


def normalize_question_type(qtype):
    """Map all type aliases to their canonical name."""
    if not qtype or not isinstance(qtype, str):
        return "mcq"
    t = qtype.strip().lower()
    if t in ["mpq", "mcq", "multiple_choice", "multiple choice", "multiple-choice"]:
        return "mcq"
    if t in ["true_false", "true false", "true/false", "true-false", "tf"]:
        return "true_false"
    if t in ["identification", "short_answer", "short answer", "short-answer"]:
        return "identification"
    if t in ["fill_blank", "fill blank", "fill-in-the-blank", "fill in the blank", "fill_in_the_blank", "fib"]:
        return "fill_blank"
    if t in ["matching", "match", "matching_pair", "matching_type"]:
        return "matching"
    return "mcq"


def safe_parse_dict(value):
    """Parse a string into a dict using JSON first, then ast.literal_eval as fallback."""
    if isinstance(value, dict):
        return value
    if not isinstance(value, str) or not value.strip():
        return {}
    try:
        parsed = json.loads(value)
        if isinstance(parsed, dict):
            return parsed
    except (json.JSONDecodeError, ValueError):
        pass
    try:
        parsed = ast.literal_eval(value)
        if isinstance(parsed, dict):
            return parsed
    except (ValueError, SyntaxError):
        pass
    return {}


def clean_extracted_text(text):
    text = str(text or '')
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def extract_text_from_file(file):
    text = ""
    try:
        filename = file.name.lower()
        if filename.endswith('.pdf'):
            reader = PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        elif filename.endswith('.docx'):
            doc = Document(file)
            for para in doc.paragraphs:
                para_text = para.text.strip()
                if para_text:
                    text += para_text + "\n"
        elif filename.endswith('.pptx'):
            prs = Presentation(file)
            for slide_index, slide in enumerate(prs.slides, start=1):
                slide_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape_text = shape.text.strip()
                        if shape_text and len(shape_text) > 1:
                            slide_lines.append(shape_text)
                if slide_lines:
                    text += f"\nSlide {slide_index}:\n"
                    text += "\n".join(slide_lines) + "\n"
        elif filename.endswith(('.png', '.jpg', '.jpeg', '.webp')):
            image = Image.open(file)
            text = pytesseract.image_to_string(image)
        else:
            content = file.read()
            if isinstance(content, bytes):
                text = content.decode('utf-8', errors='replace')
            else:
                text = str(content)
        return clean_extracted_text(text)
    except Exception as e:
        print(f"Error parsing file: {e}")
        return ""


def clean_json_text(raw_text):
    if not raw_text:
        return ""
    text = raw_text.strip()
    text = re.sub(r"^```(?:json)?", "", text, flags=re.IGNORECASE).strip()
    text = re.sub(r"```$", "", text).strip()
    return text


def extract_json_array_or_object(raw_text):
    text = clean_json_text(raw_text)
    if not text:
        return None
    try:
        return json.loads(text)
    except Exception:
        pass
    array_match = re.search(r'(\[\s*{.*}\s*\])', text, re.DOTALL)
    if array_match:
        try:
            return json.loads(array_match.group(1))
        except Exception:
            pass
    object_match = re.search(r'(\{\s*".*"\s*:.*\})', text, re.DOTALL)
    if object_match:
        try:
            return json.loads(object_match.group(1))
        except Exception:
            pass
    return None


def parse_choice_string(raw_text):
    if not isinstance(raw_text, str):
        return []
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    parsed = []
    for line in lines:
        match = re.match(r'^[A-Da-d]\s*[\).:-]\s*(.+)$', line)
        if match:
            parsed.append(match.group(1).strip())
            continue
        match = re.match(r'^([A-Da-d])\s+(.+)$', line)
        if match:
            parsed.append(match.group(2).strip())
            continue
        if '|' in line and line.count('|') == 3:
            parts = [part.strip() for part in line.split('|') if part.strip()]
            if len(parts) == 4:
                parsed = parts
                continue
        parsed.append(line)
    return parsed


def parse_matching_string(raw_text):
    if not isinstance(raw_text, str):
        return None
    lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
    left = []
    right = []
    for line in lines:
        if '->' in line:
            parts = [part.strip() for part in line.split('->', 1)]
        elif ':' in line and line.count(':') == 1:
            parts = [part.strip() for part in line.split(':', 1)]
        elif '-' in line and line.count('-') == 1:
            parts = [part.strip() for part in line.split('-', 1)]
        else:
            continue
        if len(parts) == 2:
            left.append(parts[0])
            right.append(parts[1])
    if left and right and len(left) == len(right):
        return {"left": left, "right": right}
    return None


def generate_answer_variations(answer):
    """
    Generate all acceptable answer variations including:
    - Original, lowercase, title case
    - Punctuation-stripped versions
    - FIRST WORD (e.g., "Ring Topology" → "ring")
    - FIRST TWO WORDS (e.g., "Central Processing Unit" → "central processing")
    - ACRONYM (e.g., "Central Processing Unit" → "cpu")
    """
    if not answer or len(answer.strip()) < 2:
        return []
    answer = str(answer).strip()
    variations = [answer]
    variations.append(answer.lower())
    variations.append(answer.title())

    # Punctuation-stripped version
    clean = re.sub(r'[^\w\s]', '', answer).strip()
    if clean and clean != answer:
        variations.append(clean)
        variations.append(clean.lower())

    # Whitespace-normalized version
    normalized = re.sub(r'\s+', ' ', answer).strip()
    if normalized != answer:
        variations.append(normalized)
        variations.append(normalized.lower())

    # Split into words
    words = answer.split()

    if len(words) > 1:
        # ── FIRST WORD variations ──
        # "Ring Topology" → "Ring", "ring"
        first_word = words[0]
        clean_first = re.sub(r'[^\w]', '', first_word)
        if len(clean_first) >= 2:
            variations.append(clean_first)
            variations.append(clean_first.lower())
            variations.append(clean_first.title())

        # ── FIRST TWO WORDS variations ──
        # "Central Processing Unit" → "Central Processing", "central processing"
        if len(words) >= 2:
            first_two = ' '.join(words[:2])
            clean_two = re.sub(r'[^\w\s]', '', first_two).strip()
            if len(clean_two) >= 3:
                variations.append(clean_two)
                variations.append(clean_two.lower())
                variations.append(clean_two.title())

        # ── ACRONYM variations ──
        # "Central Processing Unit" → "CPU", "cpu"
        # "Hyper Text Markup Language" → "HTML", "html"
        # Only build acronym from words that start with a letter
        acronym_letters = []
        for w in words:
            if w and w[0].isalpha():
                acronym_letters.append(w[0].upper())
        acronym = ''.join(acronym_letters)
        if len(acronym) >= 2:
            variations.append(acronym)
            variations.append(acronym.lower())
    else:
        # Single word — also strip punctuation for it
        clean_single = re.sub(r'[^\w]', '', answer).strip()
        if len(clean_single) >= 2 and clean_single != answer:
            variations.append(clean_single)
            variations.append(clean_single.lower())

    # Deduplicate preserving order
    seen = set()
    unique_variations = []
    for v in variations:
        v_clean = v.strip()
        if v_clean and len(v_clean) >= 2 and v_clean not in seen:
            seen.add(v_clean)
            unique_variations.append(v_clean)
    return unique_variations


def normalize_single_question(q, allowed_types=None):
    """Normalize a single question. Returns None if invalid."""
    if not allowed_types:
        allowed_types = ["mcq", "true_false", "identification", "fill_blank", "matching"]

    allowed_set = set(normalize_question_type(t) for t in allowed_types)

    raw_type = str(q.get("type") or q.get("question_type") or "").strip().lower()
    q_type = normalize_question_type(raw_type)

    if q_type not in allowed_set:
        return None

    if q_type == "flashcard":
        return None

    question_text = str(q.get("question", "")).strip()
    if not question_text:
        return None

    raw_choices = q.get("choices", [])

    if isinstance(raw_choices, str):
        if q_type in ["mcq", "true_false"]:
            parsed = parse_choice_string(raw_choices)
            raw_choices = parsed if len(parsed) >= 2 else [raw_choices]
        elif q_type == "matching":
            parsed_match = parse_matching_string(raw_choices)
            raw_choices = parsed_match if parsed_match else [raw_choices]
        else:
            raw_choices = [line.strip() for line in raw_choices.splitlines() if line.strip()]

    if isinstance(raw_choices, dict):
        if q_type in ["mcq", "true_false"]:
            raw_choices = raw_choices.get("options", raw_choices.get("choices", []))
        elif q_type == "matching":
            raw_choices = {"left": raw_choices.get("left", []), "right": raw_choices.get("right", [])}
        else:
            raw_choices = raw_choices.get("choices", [])

    if isinstance(raw_choices, list) and len(raw_choices) == 1 and isinstance(raw_choices[0], str):
        if q_type == "matching":
            parsed_match = parse_matching_string(raw_choices[0])
            raw_choices = parsed_match if parsed_match else raw_choices
        else:
            parsed = parse_choice_string(raw_choices[0])
            if len(parsed) >= 2:
                raw_choices = parsed

    if isinstance(raw_choices, list) and q_type == "matching":
        left_items, right_items = [], []
        for item in raw_choices:
            if isinstance(item, dict):
                if "left" in item and "right" in item:
                    left_items.append(item["left"])
                    right_items.append(item["right"])
                elif "term" in item and "definition" in item:
                    left_items.append(item["term"])
                    right_items.append(item["definition"])
            elif isinstance(item, (list, tuple)) and len(item) >= 2:
                left_items.append(item[0])
                right_items.append(item[1])
            elif isinstance(item, str):
                parsed_match = parse_matching_string(item)
                if parsed_match:
                    left_items.extend(parsed_match["left"])
                    right_items.extend(parsed_match["right"])
        raw_choices = {"left": left_items, "right": right_items}

    if isinstance(raw_choices, list) and q_type in ["mcq", "true_false"]:
        normalized_choices = []
        for choice in raw_choices:
            if isinstance(choice, dict) and "text" in choice:
                normalized_choices.append(str(choice["text"]).strip())
            else:
                normalized_choices.append(str(choice).strip())
        raw_choices = normalized_choices

    item = {
        "question": question_text,
        "type": q_type,
        "choices": raw_choices,
        "answer": q.get("answer", ""),
        "accepted_answers": q.get("accepted_answers", []),
        "explanation": str(q.get("explanation", "")).strip(),
    }

    if q_type == "true_false":
        question_text = item["question"].strip().rstrip('?.!')
        invalid_starters = ["what", "who", "where", "when", "how", "why", "which", "whom", "whose", "name", "list", "give", "identify", "describe", "explain", "define"]
        first_word = question_text.split()[0].lower() if question_text.split() else ""

        if first_word in invalid_starters:
            return None

        item["type"] = "true_false"
        item["question"] = question_text + "."
        item["choices"] = ["True", "False"]
        answer_value = str(item["answer"]).strip().lower()
        if answer_value in ["true", "t", "yes", "1"]:
            item["answer"] = "True"
        else:
            item["answer"] = "False"
        item["accepted_answers"] = [item["answer"].lower()]

    elif q_type == "mcq":
        item["type"] = "mcq"
        if not isinstance(item["choices"], list) or len(item["choices"]) < 2:
            return None
        elif len(item["choices"]) < 4:
            while len(item["choices"]) < 4:
                item["choices"].append(f"Option {chr(65 + len(item['choices']))}")
        elif len(item["choices"]) > 4:
            item["choices"] = item["choices"][:4]

        answer_map = {"a": 0, "b": 1, "c": 2, "d": 3}
        answer_key = str(item["answer"]).strip().lower()
        if answer_key in answer_map and answer_key not in [c.lower() for c in item["choices"]]:
            item["answer"] = item["choices"][answer_map[answer_key]]
        elif item["answer"] not in item["choices"]:
            found = False
            for choice in item["choices"]:
                if item["answer"].strip().lower() == choice.strip().lower():
                    item["answer"] = choice
                    found = True
                    break
            if not found:
                item["answer"] = item["choices"][0]
        if not item["accepted_answers"]:
            item["accepted_answers"] = [item["answer"]]

    elif q_type == "identification":
        item["type"] = "identification"
        item["choices"] = []
        answer = str(item["answer"]).strip()

        if not answer and isinstance(q.get("accepted_answers"), list) and q.get("accepted_answers"):
            answer = str(q.get("accepted_answers")[0]).strip()

        if not answer or len(answer) < 2:
            return None

        answer = re.sub(r'["\']', '', answer).strip()
        if not answer or len(answer) < 2:
            return None

        item["answer"] = answer
        existing = [str(a).strip() for a in item["accepted_answers"] if a and len(str(a).strip()) >= 2]
        item["accepted_answers"] = list(set(existing + generate_answer_variations(answer)))

    elif q_type == "fill_blank":
        item["type"] = "fill_blank"
        question = item["question"]
        answer = str(item["answer"]).strip()

        if not answer and isinstance(q.get("accepted_answers"), list) and q.get("accepted_answers"):
            answer = str(q.get("accepted_answers")[0]).strip()

        if not answer or len(answer) < 2:
            return None

        answer = re.sub(r'["\']', '', answer).strip()
        if not answer or len(answer) < 2:
            return None

        blank_pattern = r'_{3,}|\[BLANK\]|\{BLANK\}'
        blanks = list(re.finditer(blank_pattern, question))

        if len(blanks) > 1:
            result_parts = []
            last_end = 0
            for i, blank_match in enumerate(blanks):
                result_parts.append(question[last_end:blank_match.start()])
                if i == 0:
                    result_parts.append('________')
                last_end = blank_match.end()
            result_parts.append(question[last_end:])
            question = ''.join(result_parts)
        elif len(blanks) == 0:
            if answer and answer.lower() in question.lower():
                idx = question.lower().index(answer.lower())
                question = question[:idx] + '________' + question[idx + len(answer):]
            else:
                question = question.rstrip('.?!,;') + ': ________.'

        item["question"] = question
        item["answer"] = answer
        item["choices"] = []
        existing = [str(a).strip() for a in item["accepted_answers"] if a and len(str(a).strip()) >= 2]
        item["accepted_answers"] = list(set(existing + generate_answer_variations(answer)))

    elif q_type == "matching":
        item["type"] = "matching"
        if not isinstance(item["choices"], dict):
            item["choices"] = {"left": [], "right": []}

        parsed_answer = safe_parse_dict(item["answer"])
        item["answer"] = parsed_answer

        left = item["choices"].get("left", [])
        right = item["choices"].get("right", [])
        if not left and not right and item["answer"]:
            left = [str(k).strip() for k in item["answer"].keys()]
            right = [str(v).strip() for v in item["answer"].values()]

        if len(left) < 2:
            return None
        else:
            right_shuffled = right.copy()
            random.shuffle(right_shuffled)
            item["choices"] = {
                "left": [str(v).strip() for v in left],
                "right": [str(v).strip() for v in right_shuffled]
            }
            item["answer"] = {
                str(k).strip(): str(v).strip()
                for k, v in item["answer"].items()
            }
    else:
        return None

    return item


def normalize_questions(questions, allowed_types=None):
    if not allowed_types:
        allowed_types = ["mcq", "true_false", "identification", "fill_blank", "matching"]
    normalized = []
    for q in questions:
        item = normalize_single_question(q, allowed_types)
        if item is not None:
            normalized.append(item)
    return normalized


def get_type_example(qtype):
    qtype = normalize_question_type(qtype)
    examples = {
        "mcq": """{"question": "What does CPU stand for?", "type": "mcq", "choices": ["Central Process Unit", "Central Processing Unit", "Computer Personal Unit", "Control Processing Utility"], "answer": "Central Processing Unit", "accepted_answers": ["central processing unit", "cpu", "central processing"], "explanation": "CPU stands for Central Processing Unit."}""",

        "true_false": """{"question": "HTML is used to structure web pages.", "type": "true_false", "choices": ["True", "False"], "answer": "True", "accepted_answers": ["true"], "explanation": "HTML provides the structure of web pages."}""",

        "identification": """{"question": "What is the term for a collection of related data stored electronically?", "type": "identification", "choices": [], "answer": "Database", "accepted_answers": ["database"], "explanation": "A database stores related data electronically."}""",

        "fill_blank": """{"question": "The brain of the computer is called __________.", "type": "fill_blank", "choices": [], "answer": "Central Processing Unit", "accepted_answers": ["central processing unit", "cpu", "central processing"], "explanation": "The CPU is considered the brain of the computer."}""",

        "matching": """{"question": "Match Column A with Column B.", "type": "matching", "choices": {"left": ["HTML", "CSS", "Python", "Database"], "right": ["Structure of webpage", "Styling", "Programming Language", "Data storage"]}, "answer": {"HTML": "Structure of webpage", "CSS": "Styling", "Python": "Programming Language", "Database": "Data storage"}, "accepted_answers": [], "explanation": "Match each term to its description."}"""
    }
    return examples.get(qtype, examples["mcq"])


def get_type_rules(qtype):
    qtype = normalize_question_type(qtype)
    rules = {
        "mcq": """MCQ RULES:
- Question MUST start with "What", "Which", "Who", "Where", "When", or "How"
- Exactly 4 choices in a JSON array
- One correct answer that exactly matches one choice text
- All choices must be plausible and related to the topic""",

        "true_false": """TRUE/FALSE RULES:
- Question MUST be a DECLARATIVE STATEMENT ending with a period
- NEVER start with What/Why/How/Which/Who/Where/When
- It must be a STATEMENT that is True or False, NOT a question
- Answer must be exactly "True" or "False"
- choices must be ["True", "False"]""",

        "identification": """IDENTIFICATION RULES:
- Question MUST start with "What is the term for", "What device is", "What is", "Who is", or "Name the"
- Answer MUST be a specific term, name, or short phrase
- NEVER ask "why" or "how" or "explain" — ask for a TERM or NAME only
- Answer must be directly from or strongly related to the study content
- For multi-word answers like "Ring Topology", the answer field should contain the FULL term
- The system will automatically accept partial answers (first word, acronym) — you don't need to list them all
- choices must be empty array []
- Provide 2-4 accepted_answers with common variations (lowercase, no punctuation)""",

        "fill_blank": """FILL IN THE BLANK RULES:
- Question MUST have clear context BEFORE the blank
- Use exactly ONE blank with "________" (8 underscores)
- NEVER use multiple blanks in one question
- Answer must be a single term or short phrase
- The blank MUST replace a key term from the study content
- DO NOT use generic placeholders like "CPU" unless the content is about computers
- For multi-word answers, put the FULL term in the answer field (e.g., "Central Processing Unit")
- The system will automatically accept "CPU", "cpu", "Central Processing", "central", etc.
- choices must be empty array []
- Provide 2-4 accepted_answers with variations""",

        "matching": """MATCHING RULES:
- Minimum 3 pairs, maximum 5 pairs per question
- choices must be a JSON object with "left" and "right" arrays
- left = terms/names, right = definitions/descriptions
- Answer must be a JSON object mapping left items to right items
- CRITICAL: The answer keys MUST exactly match the left array items (same text, no extra periods or punctuation)
- Use terms and definitions FROM the study content"""
    }
    return rules.get(qtype, rules["mcq"])


def generate_questions_for_type(topic, content, difficulty, num_questions, question_type, focus_area=None):
    question_type = normalize_question_type(question_type)

    if not content:
        content = "General knowledge context."

    type_example = get_type_example(question_type)
    type_rules = get_type_rules(question_type)

    difficulty_instructions = {
        'easy': """EASY DIFFICULTY:
- Extract questions DIRECTLY from the study content provided below
- Use EXACT wording, definitions, facts, and terms from the material
- Keep questions simple and straightforward
- Answers should be directly stated in the content""",

        'medium': """MEDIUM DIFFICULTY:
- Rephrase questions from content (not copied exactly)
- Apply concepts to realistic scenarios
- Require some understanding, not just memorization""",

        'hard': """HARD DIFFICULTY:
- Create novel scenarios requiring deeper analysis
- Test understanding of relationships between concepts
- Use your general knowledge combined with the topic"""
    }

    difficulty_guidance = difficulty_instructions.get(difficulty.lower(), difficulty_instructions['medium'])

    focus_section = ""
    if focus_area and focus_area.strip():
        focus_section = f"""
╔══════════════════════════════════════════════════════════════════╗
║  MANDATORY FOCUS AREA: "{focus_area}"                              ║
║  ALL {num_questions} questions MUST be specifically about this focus area    ║
║  Do NOT generate questions about other parts of the topic          ║
║  If the content doesn't cover this focus area, use your knowledge  ║
╚══════════════════════════════════════════════════════════════════╝
"""

    system_msg = f"""You are a quiz question generator. You output ONLY valid JSON arrays.
You MUST generate EXACTLY {num_questions} questions of type "{question_type}".
{type_rules}

CRITICAL: Every question's "type" field MUST be exactly "{question_type}"."""

    if question_type == "matching":
        system_msg += "\nCRITICAL for matching: answer dict keys MUST exactly match choices.left items character-for-character."

    if difficulty.lower() == 'hard':
        content_section = f"""
TOPIC: {topic}
IMPORTANT: Use your general knowledge of the topic. Do NOT copy from any file content."""
    else:
        content_section = f"""
TOPIC: {topic}

STUDY CONTENT (use this to generate questions):
{content[:8000]}"""

    user_prompt = f"""Generate EXACTLY {num_questions} {question_type.upper()} questions.

{difficulty_guidance}
{focus_section}
{content_section}

FORMAT EXAMPLE (follow this exactly):
{type_example}

OUTPUT REQUIREMENTS:
1. Output ONLY a valid JSON array with {num_questions} objects
2. Each object MUST have "type": "{question_type}"
3. Follow the rules above exactly
4. Do NOT include any text outside the JSON array

⚠️ YOU MUST OUTPUT EXACTLY {num_questions} QUESTIONS ⚠️"""

    final_questions = []
    max_retries = 5
    retry_count = 0

    while len(final_questions) < num_questions and retry_count < max_retries:
        needed = num_questions - len(final_questions)

        if needed == num_questions:
            current_prompt = user_prompt
        else:
            current_prompt = f"""You only generated {len(final_questions)} valid questions. I need {num_questions} total.

Generate EXACTLY {needed} MORE questions of type "{question_type}".
{focus_section}
{type_rules}

FORMAT: {type_example}

Output ONLY a JSON array with {needed} objects. Each MUST have "type": "{question_type}"."""

        try:
            chat_completion = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": system_msg},
                    {"role": "user", "content": current_prompt}
                ],
                model="llama-3.1-8b-instant",
                temperature=0.4 if difficulty.lower() == 'easy' else (0.5 if difficulty.lower() == 'medium' else 0.6),
            )

            raw_text = chat_completion.choices[0].message.content
            data = extract_json_array_or_object(raw_text)

            if data is None:
                retry_count += 1
                continue

            if isinstance(data, dict):
                for key in ["quiz", "questions", "data", "items"]:
                    if key in data and isinstance(data[key], list):
                        data = data[key]
                        break
                else:
                    for value in data.values():
                        if isinstance(value, list):
                            data = value
                            break

            if isinstance(data, list):
                new_questions = normalize_questions(data, allowed_types=[question_type])

                existing_texts = {str(q.get("question", "")).strip().lower() for q in final_questions}
                for q in new_questions:
                    q_text = str(q.get("question", "")).strip().lower()
                    if q_text and q_text not in existing_texts:
                        final_questions.append(q)
                        existing_texts.add(q_text)

        except Exception as e:
            print(f"AI Generation Error (type={question_type}, retry {retry_count + 1}): {e}")

        retry_count += 1

    return final_questions[:num_questions]


def fallback_questions_for_type(topic, content, difficulty, question_type, focus_area=None):
    """Generate a single fallback question for a specific type. NEVER returns None."""
    question_type = normalize_question_type(question_type)

    target = focus_area if focus_area and focus_area.strip() else topic

    key_terms = []
    if content:
        words = re.findall(r'\b[A-Z][a-z]+(?:\s[A-Z][a-z]+)*\b', content[:2000])
        key_terms = list(set(words))[:10]

    term = key_terms[0] if key_terms else "concept"
    term2 = key_terms[1] if len(key_terms) > 1 else "principle"
    term3 = key_terms[2] if len(key_terms) > 2 else "framework"

    if question_type == "mcq":
        return {
            "question": f"What is an important aspect of {target}?",
            "type": "mcq",
            "choices": [
                f"A key principle involving {term}",
                f"An unrelated concept about {term2}",
                f"An incorrect approach to {term3}",
                f"A misconception about {target}"
            ],
            "answer": f"A key principle involving {term}",
            "accepted_answers": [f"a key principle involving {term}"],
            "explanation": f"Related to {target}."
        }

    elif question_type == "true_false":
        return {
            "question": f"{target} involves fundamental principles that can be systematically understood and applied.",
            "type": "true_false",
            "choices": ["True", "False"],
            "answer": "True",
            "accepted_answers": ["true"],
            "explanation": f"This is a declarative statement about {target}."
        }

    elif question_type == "identification":
        return {
            "question": f"What is a key term associated with {target}?",
            "type": "identification",
            "choices": [],
            "answer": term,
            "accepted_answers": generate_answer_variations(term),
            "explanation": f"A key term related to {target}."
        }

    elif question_type == "fill_blank":
        return {
            "question": f"In the context of {target}, {term} serves as a __________ concept.",
            "type": "fill_blank",
            "choices": [],
            "answer": "fundamental",
            "accepted_answers": generate_answer_variations("fundamental"),
            "explanation": f"Fill in the blank about {target}."
        }

    elif question_type == "matching":
        terms = key_terms[:3] if len(key_terms) >= 3 else [term, term2, term3]
        return {
            "question": "Match the terms with their descriptions.",
            "type": "matching",
            "choices": {
                "left": terms,
                "right": [
                    f"Related to {target}",
                    "Part of the underlying framework",
                    "A key component of the system"
                ]
            },
            "answer": {
                terms[0]: f"Related to {target}",
                terms[1]: "Part of the underlying framework",
                terms[2]: "A key component of the system"
            },
            "accepted_answers": [],
            "explanation": "Matching question."
        }

    return {
        "question": f"What is an important aspect of {target}?",
        "type": "mcq",
        "choices": [
            f"A key principle involving {term}",
            f"An unrelated concept",
            f"An incorrect approach",
            f"A misconception"
        ],
        "answer": f"A key principle involving {term}",
        "accepted_answers": [f"a key principle involving {term}"],
        "explanation": f"Related to {target}."
    }


def generate_quiz_questions(topic, content, difficulty, num_questions, question_types, focus_area=None):
    if not content:
        content = "General knowledge context."

    question_types = [normalize_question_type(q) for q in question_types]
    question_types = [q for q in question_types if q not in ['flashcard']]

    if not question_types:
        question_types = ["mcq"]

    question_types = list(dict.fromkeys(question_types))

    num_types = len(question_types)
    base_count = num_questions // num_types
    remainder = num_questions % num_types

    type_counts = {}
    for i, qtype in enumerate(question_types):
        type_counts[qtype] = base_count + (1 if i < remainder else 1)

    total = sum(type_counts.values())
    while total > num_questions:
        max_type = max(type_counts, key=type_counts.get)
        type_counts[max_type] -= 1
        total = sum(type_counts.values())

    for qtype in question_types:
        if type_counts.get(qtype, 0) < 1:
            type_counts[qtype] = 1

    total = sum(type_counts.values())
    if total != num_questions:
        diff = num_questions - total
        if diff > 0:
            for qtype in question_types:
                type_counts[qtype] += 1
                diff -= 1
                if diff == 0:
                    break
        elif diff < 0:
            while diff < 0:
                max_type = max(type_counts, key=type_counts.get)
                if type_counts[max_type] > 1:
                    type_counts[max_type] -= 1
                    diff += 1
                else:
                    break

    print(f"Question distribution: {type_counts}")

    all_questions = []

    for qtype, count in type_counts.items():
        if count <= 0:
            continue

        print(f"Generating {count} {qtype} questions...")

        type_questions = generate_questions_for_type(
            topic=topic,
            content=content,
            difficulty=difficulty,
            num_questions=count,
            question_type=qtype,
            focus_area=focus_area
        )

        print(f"Got {len(type_questions)} valid {qtype} questions")

        if len(type_questions) < count:
            needed = count - len(type_questions)
            print(f"Adding {needed} fallback {qtype} questions")
            for i in range(needed):
                fallback = fallback_questions_for_type(topic, content, difficulty, qtype, focus_area)
                if fallback:
                    fallback_text = fallback["question"].strip().lower()
                    existing_texts = {q["question"].strip().lower() for q in type_questions}
                    if fallback_text not in existing_texts:
                        type_questions.append(fallback)
                    else:
                        fallback["question"] = fallback["question"].rstrip('.') + f" (#{i+1})."
                        type_questions.append(fallback)

        all_questions.extend(type_questions[:count])

    if len(all_questions) < num_questions:
        needed = num_questions - len(all_questions)
        print(f"FINAL PAD: Adding {needed} extra fallback questions")
        for qtype in question_types:
            for i in range(needed):
                fallback = fallback_questions_for_type(topic, content, difficulty, qtype, focus_area)
                if fallback:
                    all_questions.append(fallback)
                    needed -= 1
                    if needed <= 0:
                        break
            if needed <= 0:
                break

    random.shuffle(all_questions)

    result = all_questions[:num_questions]
    print(f"Final question count: {len(result)}")

    return result